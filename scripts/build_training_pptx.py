#!/usr/bin/env python3
"""Build trainer PPTX: title + infographics + closing (no Marp slides)."""
from __future__ import annotations

import sys
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
INFO_DIR = ROOT / "trainer" / "infographics"
OUT_PPTX = ROOT / "trainer" / "training-session-20260605.pptx"
OUT_PPTX_ALT = ROOT / "trainer" / "training-session-20260605-new.pptx"
FEEDBACK_URL = "https://forms.gle/efbJ3VZwB8C5q4Pa7"
QR_PATH = ROOT / "trainer" / "feedback-form-qrcode.png"

INFO_SLIDES: list[tuple[str, str]] = [
    ("fig-02-why-cursor-hk.png", "§0 點解學 Cursor"),
    ("fig-01-cursor-interface.png", "§1 Cursor 介面"),
    ("fig-06-chatbot-to-agent.png", "§2 AI 2026 變革"),
    ("fig-03-five-concepts.png", "五個核心概念"),
    ("fig-05-90min-timeline.png", "90 分鐘路線圖"),
    ("fig-07-activity1-workflow.png", "活動一 Workflow"),
    ("fig-04-boss-vs-intern.png", "Vibe Coding · 老闘 vs 實習文員"),
    ("fig-08-activity2-files.png", "活動二 本機執檔"),
    ("fig-09-activity3-marp.png", "活動三 靜態網站"),
]

CLOSING = ("fig-10-closing-summary.png", "總結 · 工作流總設計師")


def generate_qr_code(url: str, path: Path) -> Path:
    qr = qrcode.QRCode(version=1, box_size=12, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(path))
    return path


def add_feedback_slide(prs: Presentation, url: str = FEEDBACK_URL) -> None:
    qr_path = generate_qr_code(url, QR_PATH)
    slide = prs.slides.add_slide(blank_layout(prs))

    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "課後回饋問卷"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "請掃描 QR Code 或開啟下方連結 · 約 3–5 分鐘"
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
    sp.alignment = PP_ALIGN.CENTER

    qr_size = Inches(3.2)
    left = (prs.slide_width - qr_size) / 2
    slide.shapes.add_picture(str(qr_path), left, Inches(2.3), width=qr_size, height=qr_size)

    link_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.75), Inches(10.3), Inches(0.8))
    lp = link_box.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    run = lp.add_run()
    run.text = url
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    run.hyperlink.address = url

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5))
    fp = foot.text_frame.paragraphs[0]
    fp.text = "感謝您的寶貴意見 · Carmel Holy Word Secondary School"
    fp.font.size = Pt(14)
    fp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    fp.alignment = PP_ALIGN.CENTER


def blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name.lower() in {"blank", "空白"}:
            return layout
    return prs.slide_layouts[6]


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CHW 教職員 Cursor 培訓"
    p.font.size = Inches(0.45)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "2026-06-05 · Agentic Workflow"
    p2.font.size = Inches(0.28)
    p3 = tf.add_paragraph()
    p3.text = "迦密聖道中學 Carmel Holy Word Secondary School"
    p3.font.size = Inches(0.22)


def add_image_slide(prs: Presentation, image_path: Path, note: str = "") -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    slide.shapes.add_picture(
        str(image_path),
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )
    if note:
        box = slide.shapes.add_textbox(
            Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35)
        )
        tf = box.text_frame
        tf.paragraphs[0].text = note
        tf.paragraphs[0].font.size = Inches(0.14)


def main() -> int:
    if not INFO_DIR.exists():
        print(f"Missing infographics dir: {INFO_DIR}", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    for filename, note in INFO_SLIDES:
        path = INFO_DIR / filename
        if not path.exists():
            print(f"Missing: {path}", file=sys.stderr)
            return 1
        add_image_slide(prs, path, note)
        print(f"  + infographic: {filename}")

    closing_path = INFO_DIR / CLOSING[0]
    if closing_path.exists():
        add_image_slide(prs, closing_path, CLOSING[1])
        print(f"  + infographic: {CLOSING[0]}")

    add_feedback_slide(prs)
    print(f"  + feedback slide: {FEEDBACK_URL}")

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    saved: list[Path] = []
    for target in [OUT_PPTX, OUT_PPTX_ALT]:
        try:
            prs.save(str(target))
            saved.append(target)
            print(f"Saved: {target} ({len(prs.slides)} slides)")
            break
        except PermissionError:
            print(f"  ! locked, skip: {target}", file=sys.stderr)
    if not saved:
        fallback = OUT_PPTX_ALT
        prs.save(str(fallback))
        saved.append(fallback)
        print(f"Saved: {fallback} ({len(prs.slides)} slides)")

    return 0


if __name__ == "__main__":
    sys.exit(main())
